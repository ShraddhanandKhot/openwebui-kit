#!/usr/bin/env python3
"""
Filtration Layer - Complexity Detector for Document Processing Pipeline

This module inspects files BEFORE conversion to determine:
- Simple files: Direct to markdown conversion (fast path)
- Complex files: Need pre-processing/extraction (slow path)

Complexity factors:
1. Scanned/image-based PDFs (need OCR)
2. Embedded images (need extraction + description)
3. Complex tables (need special handling)
4. Comments/annotations (need extraction)
5. Forms/fields (need extraction)
6. Multi-column layouts (need layout analysis)
7. Equations/formulas (need special handling)
8. Watermarks/backgrounds (may interfere)
"""

import os
import json
import sys
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any
from dataclasses import dataclass, field
from enum import Enum


class ComplexityLevel(Enum):
    SIMPLE = "simple"      # Direct to markdown
    MODERATE = "moderate"  # Some pre-processing needed
    COMPLEX = "complex"    # Full extraction pipeline needed


@dataclass
class FileComplexityReport:
    """Report on file complexity analysis."""
    file_path: str
    file_type: str
    file_size_bytes: int
    complexity_level: ComplexityLevel
    complexity_score: int  # 0-100
    factors: List[str] = field(default_factory=list)
    recommendations: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self) -> Dict:
        return {
            "file_path": self.file_path,
            "file_type": self.file_type,
            "file_size_bytes": self.file_size_bytes,
            "complexity_level": self.complexity_level.value,
            "complexity_score": self.complexity_score,
            "factors": self.factors,
            "recommendations": self.recommendations,
            "metadata": self.metadata
        }


class FiltrationLayer:
    """
    Filtration Layer - Analyzes files before markdown conversion.
    
    Determines the optimal processing path based on content complexity.
    """
    
    # File extensions that are typically simple (text-based)
    SIMPLE_EXTENSIONS = {
        '.txt', '.md', '.markdown', '.rst', '.csv', '.tsv',
        '.json', '.xml', '.yaml', '.yml', '.html', '.htm',
        '.log', '.ini', '.cfg', '.conf', '.py', '.js', '.ts',
        '.java', '.cpp', '.c', '.h', '.go', '.rs', '.rb',
        '.php', '.swift', '.kt', '.scala', '.sh', '.bash',
        '.sql', '.css', '.scss', '.less'
    }
    
    # Extensions that often contain complex elements
    COMPLEX_EXTENSIONS = {
        '.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls',
        '.odt', '.ods', '.odp', '.epub', '.rtf'
    }
    
    # Image extensions (always complex - need OCR/description)
    IMAGE_EXTENSIONS = {
        '.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp',
        '.gif', '.webp', '.heic', '.heif'
    }
    
    def __init__(self):
        self._init_analyzers()
    
    def _init_analyzers(self):
        """Initialize available analyzers."""
        self.has_pymupdf = False
        self.has_unstructured = False
        self.has_pytesseract = False
        
        try:
            import pymupdf
            self.has_pymupdf = True
        except ImportError:
            pass
            
        try:
            from unstructured.partition.auto import partition
            self.has_unstructured = True
        except ImportError:
            pass
            
        try:
            import pytesseract
            self.has_pytesseract = True
        except ImportError:
            pass
    
    def analyze_file(self, file_path: str) -> FileComplexityReport:
        """
        Main entry point - analyze a file and return complexity report.
        """
        path = Path(file_path)
        
        if not path.exists():
            return FileComplexityReport(
                file_path=file_path,
                file_type="unknown",
                file_size_bytes=0,
                complexity_level=ComplexityLevel.COMPLEX,
                complexity_score=100,
                factors=["file_not_found"],
                recommendations=["File does not exist"]
            )
        
        file_size = path.stat().st_size
        file_ext = path.suffix.lower()
        file_type = self._detect_file_type(file_path, file_ext)
        
        # Quick path for known simple extensions
        if file_ext in self.SIMPLE_EXTENSIONS:
            return FileComplexityReport(
                file_path=file_path,
                file_type=file_type,
                file_size_bytes=file_size,
                complexity_level=ComplexityLevel.SIMPLE,
                complexity_score=10,
                factors=["text_based_extension"],
                recommendations=["Direct markdown conversion"],
                metadata={"extension": file_ext}
            )
        
        # Images always need OCR/vision
        if file_ext in self.IMAGE_EXTENSIONS:
            return FileComplexityReport(
                file_path=file_path,
                file_type=file_type,
                file_size_bytes=file_size,
                complexity_level=ComplexityLevel.COMPLEX,
                complexity_score=85,
                factors=["image_file_needs_ocr"],
                recommendations=["OCR required", "Use vision model for description"],
                metadata={"extension": file_ext}
            )
        
        # Deep analysis for complex document types
        return self._deep_analyze(file_path, file_type, file_size, file_ext)
    
    def _detect_file_type(self, file_path: str, ext: str) -> str:
        """Detect actual file type using magic bytes."""
        try:
            import magic
            mime = magic.from_file(file_path, mime=True)
            return mime
        except ImportError:
            # Fallback to extension
            type_map = {
                '.pdf': 'application/pdf',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.doc': 'application/msword',
                '.ppt': 'application/vnd.ms-powerpoint',
                '.xls': 'application/vnd.ms-excel',
                '.odt': 'application/vnd.oasis.opendocument.text',
                '.ods': 'application/vnd.oasis.opendocument.spreadsheet',
                '.odp': 'application/vnd.oasis.opendocument.presentation',
                '.epub': 'application/epub+zip',
                '.rtf': 'application/rtf',
            }
            return type_map.get(ext, 'application/octet-stream')
    
    def _deep_analyze(self, file_path: str, file_type: str, file_size: int, ext: str) -> FileComplexityReport:
        """Deep content analysis for complex document types."""
        factors = []
        recommendations = []
        metadata = {"extension": ext, "file_type": file_type}
        score = 20  # Base score for non-text files
        
        if 'pdf' in file_type or ext == '.pdf':
            return self._analyze_pdf(file_path, file_size, factors, recommendations, metadata, score)
        elif 'word' in file_type or ext in ['.docx', '.doc']:
            return self._analyze_docx(file_path, file_size, factors, recommendations, metadata, score)
        elif 'presentation' in file_type or ext in ['.pptx', '.ppt']:
            return self._analyze_pptx(file_path, file_size, factors, recommendations, metadata, score)
        elif 'spreadsheet' in file_type or ext in ['.xlsx', '.xls']:
            return self._analyze_xlsx(file_path, file_size, factors, recommendations, metadata, score)
        elif 'opendocument' in file_type or ext in ['.odt', '.ods', '.odp']:
            return self._analyze_opendoc(file_path, file_size, factors, recommendations, metadata, score)
        elif 'epub' in file_type or ext == '.epub':
            return self._analyze_epub(file_path, file_size, factors, recommendations, metadata, score)
        elif 'rtf' in file_type or ext == '.rtf':
            return self._analyze_rtf(file_path, file_size, factors, recommendations, metadata, score)
        else:
            return FileComplexityReport(
                file_path=file_path,
                file_type=file_type,
                file_size_bytes=file_size,
                complexity_level=ComplexityLevel.MODERATE,
                complexity_score=40,
                factors=["unknown_document_type"],
                recommendations=["Try unstructured partition"],
                metadata=metadata
            )
    
    def _analyze_pdf(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze PDF for complexity factors."""
        score = base_score
        
        if not self.has_pymupdf:
            factors.append("pymupdf_not_available")
            recommendations.append("Install pymupdf for better analysis")
            score += 20
        else:
            import pymupdf
            doc = pymupdf.open(file_path)
            metadata["page_count"] = len(doc)
            
            # Check if scanned (no text, only images)
            total_text_chars = 0
            total_images = 0
            has_tables = False
            has_forms = False
            has_annotations = False
            has_links = False
            
            for page_num, page in enumerate(doc):
                text = page.get_text()
                total_text_chars += len(text.strip())
                
                # Count images
                images = page.get_images(full=True)
                total_images += len(images)
                
                # Check for tables (using find_tables if available)
                try:
                    tables = page.find_tables()
                    if tables.tables:
                        has_tables = True
                        factors.append(f"table_on_page_{page_num+1}")
                except:
                    pass
                
                # Check for form fields
                try:
                    widgets = list(page.widgets() or [])
                    if widgets:
                        has_forms = True
                        factors.append(f"form_fields_on_page_{page_num+1}")
                except:
                    pass

                # Check for annotations/comments
                try:
                    annots = list(page.annots() or [])
                    if annots:
                        has_annotations = True
                        factors.append(f"annotations_on_page_{page_num+1}")
                except:
                    pass
                
                # Check for links
                links = page.get_links()
                if links:
                    has_links = True
                    factors.append(f"links_on_page_{page_num+1}")
            
            metadata["total_text_chars"] = total_text_chars
            metadata["total_images"] = total_images
            metadata["has_tables"] = has_tables
            metadata["has_forms"] = has_forms
            metadata["has_annotations"] = has_annotations
            metadata["has_links"] = has_links
            
            # Determine if scanned (very little text, many images)
            if total_text_chars < 100 and total_images > 0:
                factors.append("likely_scanned_pdf")
                score += 50
                recommendations.append("OCR required - appears to be scanned document")
            elif total_text_chars < 500 and total_images > len(doc) * 0.5:
                factors.append("mostly_images_low_text")
                score += 30
                recommendations.append("May need OCR for image content")
            
            if has_tables:
                score += 15
                recommendations.append("Table extraction needed")
            
            if has_forms:
                score += 10
                recommendations.append("Form field extraction needed")
            
            if has_annotations:
                score += 10
                recommendations.append("Annotation/comment extraction needed")
            
            if total_images > 0:
                score += min(total_images * 2, 20)
                recommendations.append(f"Extract {total_images} embedded images")
        
        # Determine final complexity level
        if score <= 30:
            level = ComplexityLevel.SIMPLE
        elif score <= 60:
            level = ComplexityLevel.MODERATE
        else:
            level = ComplexityLevel.COMPLEX
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/pdf",
            file_size_bytes=file_size,
            complexity_level=level,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )
    
    def _analyze_docx(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze DOCX for complexity factors."""
        score = base_score
        
        try:
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(file_path)
            
            # Count paragraphs, tables, images
            para_count = len(doc.paragraphs)
            table_count = len(doc.tables)
            image_count = 0
            
            # Count images in document
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1
            
            # Check for comments
            has_comments = False
            try:
                from docx.oxml.comments import CT_Comment
                if doc.part.comments_part:
                    has_comments = True
            except:
                pass
            
            metadata["paragraph_count"] = para_count
            metadata["table_count"] = table_count
            metadata["image_count"] = image_count
            metadata["has_comments"] = has_comments
            
            if table_count > 0:
                factors.append(f"docx_tables_{table_count}")
                score += min(table_count * 5, 20)
                recommendations.append(f"Extract {table_count} tables")
            
            if image_count > 0:
                factors.append(f"docx_images_{image_count}")
                score += min(image_count * 3, 20)
                recommendations.append(f"Extract {image_count} embedded images")
            
            if has_comments:
                factors.append("docx_comments")
                score += 10
                recommendations.append("Extract comments/track changes")
            
            # Check for complex formatting
            style_count = len(set(p.style.name for p in doc.paragraphs if p.style))
            if style_count > 10:
                factors.append("complex_styling")
                score += 10
            
        except Exception as e:
            factors.append(f"docx_analysis_error: {e}")
            score += 15
            recommendations.append("Could not fully analyze DOCX structure")
        
        if score <= 30:
            level = ComplexityLevel.SIMPLE
        elif score <= 60:
            level = ComplexityLevel.MODERATE
        else:
            level = ComplexityLevel.COMPLEX
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            file_size_bytes=file_size,
            complexity_level=level,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )
    
    def _analyze_pptx(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze PPTX for complexity factors."""
        score = base_score + 20  # Presentations are inherently more complex
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            slide_count = len(prs.slides)
            total_shapes = 0
            table_count = 0
            image_count = 0
            has_notes = False
            
            for slide in prs.slides:
                total_shapes += len(slide.shapes)
                for shape in slide.shapes:
                    if shape.has_table:
                        table_count += 1
                    if hasattr(shape, 'image'):
                        image_count += 1
                
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        has_notes = True
            
            metadata["slide_count"] = slide_count
            metadata["total_shapes"] = total_shapes
            metadata["table_count"] = table_count
            metadata["image_count"] = image_count
            metadata["has_notes"] = has_notes
            
            if table_count > 0:
                factors.append(f"pptx_tables_{table_count}")
                score += min(table_count * 5, 15)
            
            if image_count > 0:
                factors.append(f"pptx_images_{image_count}")
                score += min(image_count * 3, 20)
            
            if has_notes:
                factors.append("pptx_speaker_notes")
                score += 5
                recommendations.append("Extract speaker notes")
                
        except Exception as e:
            factors.append(f"pptx_analysis_error: {e}")
            score += 15
        
        if score <= 35:
            level = ComplexityLevel.SIMPLE
        elif score <= 65:
            level = ComplexityLevel.MODERATE
        else:
            level = ComplexityLevel.COMPLEX
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            file_size_bytes=file_size,
            complexity_level=level,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )
    
    def _analyze_xlsx(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze XLSX for complexity factors."""
        score = base_score + 10
        
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            
            sheet_names = wb.sheetnames
            total_rows = 0
            total_cols = 0
            has_formulas = False
            has_merged = False
            has_charts = False
            has_images = False
            
            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                total_rows += ws.max_row
                total_cols = max(total_cols, ws.max_column)
                
                # Check for merged cells (read_only worksheets may lack this attr)
                merged = getattr(ws, "merged_cells", None)
                if merged is not None and getattr(merged, "ranges", None):
                    has_merged = True
                
                # Check for formulas
                for row in ws.iter_rows(max_row=min(100, ws.max_row), values_only=False):
                    for cell in row:
                        if cell.value and isinstance(cell.value, str) and cell.value.startswith('='):
                            has_formulas = True
                            break
            
            # Check for charts/images (requires full load)
            try:
                wb_full = load_workbook(file_path, read_only=False)
                for ws in wb_full.worksheets:
                    if ws._charts:
                        has_charts = True
                    for rel in ws._rels.values():
                        if "image" in str(rel.target_ref):
                            has_images = True
            except:
                pass
            
            metadata["sheet_count"] = len(sheet_names)
            metadata["sheet_names"] = sheet_names
            metadata["total_rows"] = total_rows
            metadata["max_columns"] = total_cols
            metadata["has_formulas"] = has_formulas
            metadata["has_merged_cells"] = has_merged
            metadata["has_charts"] = has_charts
            metadata["has_images"] = has_images
            
            if has_merged:
                factors.append("xlsx_merged_cells")
                score += 10
            
            if has_formulas:
                factors.append("xlsx_formulas")
                score += 10
            
            if has_charts:
                factors.append("xlsx_charts")
                score += 15
                recommendations.append("Extract chart data")
            
            if has_images:
                factors.append("xlsx_images")
                score += 10
            
            if total_rows > 10000:
                factors.append("large_dataset")
                score += 10
                recommendations.append("Large dataset - consider sampling")
                
        except Exception as e:
            factors.append(f"xlsx_analysis_error: {e}")
            score += 15
        
        if score <= 30:
            level = ComplexityLevel.SIMPLE
        elif score <= 60:
            level = ComplexityLevel.MODERATE
        else:
            level = ComplexityLevel.COMPLEX
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            file_size_bytes=file_size,
            complexity_level=level,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )
    
    def _analyze_opendoc(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze OpenDocument formats."""
        score = base_score + 15
        factors.append("opendocument_format")
        recommendations.append("Use unstructured for OpenDocument parsing")
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/vnd.oasis.opendocument",
            file_size_bytes=file_size,
            complexity_level=ComplexityLevel.MODERATE,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )
    
    def _analyze_epub(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze EPUB."""
        score = base_score + 10
        factors.append("epub_format")
        recommendations.append("Use unstructured for EPUB parsing")
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/epub+zip",
            file_size_bytes=file_size,
            complexity_level=ComplexityLevel.MODERATE,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )
    
    def _analyze_rtf(self, file_path: str, file_size: int, factors: List, recommendations: List, metadata: Dict, base_score: int) -> FileComplexityReport:
        """Analyze RTF."""
        score = base_score + 5
        factors.append("rtf_format")
        recommendations.append("Use unstructured for RTF parsing")
        
        return FileComplexityReport(
            file_path=file_path,
            file_type="application/rtf",
            file_size_bytes=file_size,
            complexity_level=ComplexityLevel.SIMPLE,
            complexity_score=min(score, 100),
            factors=factors,
            recommendations=recommendations,
            metadata=metadata
        )


def analyze_file_complexity(file_path: str) -> FileComplexityReport:
    """Convenience function to analyze a single file."""
    filtration = FiltrationLayer()
    return filtration.analyze_file(file_path)


def batch_analyze(file_paths: List[str]) -> List[FileComplexityReport]:
    """Analyze multiple files."""
    filtration = FiltrationLayer()
    return [filtration.analyze_file(fp) for fp in file_paths]


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python filtration_layer.py <file_path> [--json]")
        sys.exit(1)
    
    file_path = sys.argv[1]
    as_json = "--json" in sys.argv
    
    report = analyze_file_complexity(file_path)
    
    if as_json:
        print(json.dumps(report.to_dict(), indent=2))
    else:
        print(f"File: {report.file_path}")
        print(f"Type: {report.file_type}")
        print(f"Size: {report.file_size_bytes} bytes")
        print(f"Complexity: {report.complexity_level.value.upper()} (score: {report.complexity_score}/100)")
        print(f"Factors: {', '.join(report.factors) if report.factors else 'none'}")
        print(f"Recommendations:")
        for rec in report.recommendations:
            print(f"  - {rec}")