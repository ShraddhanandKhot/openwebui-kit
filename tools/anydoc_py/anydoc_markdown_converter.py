#!/usr/bin/env python3
"""
AnyDoc Markdown Converter - Universal Document to Markdown Pipeline

This module converts ANY document to clean, meaningful Markdown using:
1. Filtration Layer - determines optimal processing path
2. Multiple conversion backends (unstructured, pymupdf4llm, marker-pdf, python-docx, etc.)
3. Post-processing to clean and structure the output

Supports: PDF, DOCX, PPTX, XLSX, ODT, ODS, ODP, EPUB, RTF, HTML, Images, Text formats
"""

import os
import json
import sys
import tempfile
import base64
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, BinaryIO
from dataclasses import dataclass, field
from io import BytesIO

# Import our filtration layer
from filtration_layer import FiltrationLayer, FileComplexityReport, ComplexityLevel, analyze_file_complexity


@dataclass
class ConversionResult:
    """Result of a document-to-markdown conversion."""
    success: bool
    markdown: str = ""
    error: str = ""
    complexity_report: Optional[FileComplexityReport] = None
    extraction_method: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    images_extracted: List[Dict] = field(default_factory=list)
    tables_extracted: List[Dict] = field(default_factory=list)
    
    def to_dict(self) -> Dict:
        result = {
            "success": self.success,
            "markdown": self.markdown,
            "error": self.error,
            "extraction_method": self.extraction_method,
            "metadata": self.metadata,
            "images_extracted": self.images_extracted,
            "tables_extracted": self.tables_extracted,
        }
        if self.complexity_report:
            result["complexity_report"] = self.complexity_report.to_dict()
        return result


class AnyDocMarkdownConverter:
    """
    Universal document to Markdown converter with intelligent routing.
    
    Uses Filtration Layer to determine the best conversion strategy:
    - Simple files: Fast path (pymupdf4llm, unstructured, python-docx)
    - Moderate files: Enhanced extraction with table/image handling
    - Complex files: Full pipeline (OCR, layout analysis, image description)
    """
    
    def __init__(self, enable_ocr: bool = False, enable_vision: bool = False):
        self.filtration = FiltrationLayer()
        self.enable_ocr = enable_ocr
        self.enable_vision = enable_vision
        self._init_backends()
    
    def _init_backends(self):
        """Check available conversion backends."""
        self.backends = {
            "pymupdf4llm": False,
            "unstructured": False,
            "marker_pdf": False,
            "python_docx": False,
            "python_pptx": False,
            "openpyxl": False,
            "pytesseract": False,
        }
        
        try:
            import pymupdf4llm
            self.backends["pymupdf4llm"] = True
        except ImportError:
            pass
        
        try:
            from unstructured.partition.auto import partition
            self.backends["unstructured"] = True
        except ImportError:
            pass
        
        try:
            from marker.converters.pdf import PdfConverter
            self.backends["marker_pdf"] = True
        except ImportError:
            pass
        
        try:
            from docx import Document
            self.backends["python_docx"] = True
        except ImportError:
            pass
        
        try:
            from pptx import Presentation
            self.backends["python_pptx"] = True
        except ImportError:
            pass
        
        try:
            import openpyxl
            self.backends["openpyxl"] = True
        except ImportError:
            pass
        
        try:
            import pytesseract
            self.backends["pytesseract"] = True
        except ImportError:
            pass
    
    def convert(self, file_path: str, **options) -> ConversionResult:
        """
        Main conversion entry point.
        
        Args:
            file_path: Path to the document file
            **options: Additional options (extract_images, extract_tables, etc.)
        
        Returns:
            ConversionResult with markdown and metadata
        """
        # Step 1: Analyze complexity
        complexity_report = analyze_file_complexity(file_path)
        
        # Step 2: Route to appropriate converter
        if complexity_report.complexity_level == ComplexityLevel.SIMPLE:
            return self._convert_simple(file_path, complexity_report, **options)
        elif complexity_report.complexity_level == ComplexityLevel.MODERATE:
            return self._convert_moderate(file_path, complexity_report, **options)
        else:
            return self._convert_complex(file_path, complexity_report, **options)
    
    def convert_bytes(self, file_bytes: bytes, filename: str, **options) -> ConversionResult:
        """Convert from bytes (for uploaded files)."""
        with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        
        try:
            result = self.convert(tmp_path, **options)
            return result
        finally:
            try:
                os.unlink(tmp_path)
            except:
                pass
    
    def _convert_simple(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Fast path for simple files."""
        ext = Path(file_path).suffix.lower()
        
        # Try pymupdf4llm first for PDFs (best quality)
        if ext == '.pdf' and self.backends["pymupdf4llm"]:
            return self._convert_pdf_pymupdf4llm(file_path, report, **options)
        
        # Try unstructured for most formats
        if self.backends["unstructured"]:
            return self._convert_unstructured(file_path, report, **options)
        
        # Fallback to format-specific converters
        if ext == '.pdf' and self.backends["pymupdf4llm"]:
            return self._convert_pdf_pymupdf4llm(file_path, report, **options)
        elif ext in ['.docx', '.doc'] and self.backends["python_docx"]:
            return self._convert_docx_native(file_path, report, **options)
        elif ext in ['.pptx', '.ppt'] and self.backends["python_pptx"]:
            return self._convert_pptx_native(file_path, report, **options)
        elif ext in ['.xlsx', '.xls'] and self.backends["openpyxl"]:
            return self._convert_xlsx_native(file_path, report, **options)
        elif ext in ['.html', '.htm'] and self.backends["unstructured"]:
            return self._convert_unstructured(file_path, report, **options)
        elif ext in ['.md', '.markdown', '.txt', '.rst']:
            return self._convert_text_file(file_path, report, **options)
        
        return ConversionResult(
            success=False,
            error="No suitable converter available for this file type",
            complexity_report=report,
            extraction_method="none"
        )
    
    def _convert_moderate(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Enhanced conversion for moderate complexity files."""
        ext = Path(file_path).suffix.lower()
        
        # For PDFs with tables/images, try enhanced extraction
        if ext == '.pdf':
            if self.backends["pymupdf4llm"]:
                opts = dict(options)
                opts.pop("extract_images", None)
                result = self._convert_pdf_pymupdf4llm(file_path, report, extract_images=options.get("extract_images", True), **opts)
                if result.success and (report.metadata.get("has_tables") or report.metadata.get("total_images", 0) > 0):
                    # Enhance with table extraction
                    enhanced = self._enhance_with_tables(file_path, result, report)
                    if enhanced:
                        return enhanced
                return result
        
        # For DOCX with tables/images
        if ext in ['.docx', '.doc'] and self.backends["python_docx"]:
            return self._convert_docx_enhanced(file_path, report, **options)
        
        # For PPTX
        if ext in ['.pptx', '.ppt'] and self.backends["python_pptx"]:
            return self._convert_pptx_enhanced(file_path, report, **options)
        
        # For XLSX
        if ext in ['.xlsx', '.xls'] and self.backends["openpyxl"]:
            return self._convert_xlsx_enhanced(file_path, report, **options)
        
        # Fallback to unstructured
        if self.backends["unstructured"]:
            return self._convert_unstructured(file_path, report, **options)
        
        return self._convert_simple(file_path, report, **options)
    
    def _convert_complex(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Full pipeline for complex files (scanned, heavy images, etc.)."""
        ext = Path(file_path).suffix.lower()
        
        # For scanned PDFs - try marker-pdf if available (best OCR)
        if ext == '.pdf' and self.backends["marker_pdf"]:
            return self._convert_pdf_marker(file_path, report, **options)
        
        # For images - need OCR
        if ext in ['.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.webp']:
            return self._convert_image_ocr(file_path, report, **options)
        
        # For PDFs without marker - try pymupdf4llm with OCR hint
        if ext == '.pdf' and self.backends["pymupdf4llm"]:
            result = self._convert_pdf_pymupdf4llm(file_path, report, use_ocr=True, **options)
            # pymupdf4llm's bundled OCR can silently return EMPTY on some
            # image-only PDFs -> fall back to per-page tesseract OCR.
            if result.success and len((result.markdown or "").strip()) >= 40:
                return result
            ocr_result = self._convert_pdf_pytesseract_ocr(file_path, report, **options)
            if ocr_result.success and len((ocr_result.markdown or "").strip()) >= 40:
                return ocr_result
            return result
        
        # Fallback to unstructured with OCR
        if self.backends["unstructured"] and self.enable_ocr:
            return self._convert_unstructured_ocr(file_path, report, **options)
        
        return ConversionResult(
            success=False,
            error="Complex file requires OCR/marker-pdf which is not available",
            complexity_report=report,
            extraction_method="none"
        )
    
    # === PDF Converters ===
    
    def _convert_pdf_pymupdf4llm(self, file_path: str, report: FileComplexityReport, 
                                  extract_images: bool = False, use_ocr: bool = False, **options) -> ConversionResult:
        """Convert PDF using pymupdf4llm (high quality, fast)."""
        try:
            import pymupdf4llm
            import pymupdf
            
            # Configure options
            md_options = {
                "page_chunks": False,
                "write_images": extract_images,
                "image_path": options.get("image_output_dir", "./extracted_images"),
                "image_format": "png",
                "dpi": options.get("dpi", 150),
            }
            
            if use_ocr:
                md_options["ocr"] = True
            
            markdown = pymupdf4llm.to_markdown(file_path, **md_options)
            
            # Extract images if requested
            images = []
            if extract_images:
                images = self._extract_pdf_images(file_path, md_options.get("image_path", "./extracted_images"))
            
            # Extract tables
            tables = self._extract_pdf_tables(file_path)
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="pymupdf4llm",
                metadata={"engine": "pymupdf4llm", "use_ocr": use_ocr},
                images_extracted=images,
                tables_extracted=tables
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"pymupdf4llm conversion failed: {e}",
                complexity_report=report,
                extraction_method="pymupdf4llm"
            )
    
    def _convert_pdf_pytesseract_ocr(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """OCR an image-only PDF page-by-page with tesseract (via pytesseract).

        Renders each page to PNG with pymupdf, runs tesseract OCR, builds markdown.
        Used as a fallback when pymupdf4llm's bundled OCR returns empty.
        """
        try:
            import pymupdf
            import pytesseract
            from PIL import Image

            doc = pymupdf.open(file_path)
            md_parts = []
            with tempfile.TemporaryDirectory() as tmpdir:
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=200)
                    img_path = os.path.join(tmpdir, f"page_{i+1}.png")
                    pix.save(img_path)
                    text = pytesseract.image_to_string(Image.open(img_path))
                    text = text.strip()
                    if text:
                        md_parts.append(f"## Page {i+1}\n\n{text}")
                    else:
                        md_parts.append(f"## Page {i+1}\n\n_(no readable text on this page)_")
            doc.close()

            markdown = "\n\n".join(md_parts)
            if not markdown.strip() or not any("_(no readable" not in p for p in md_parts):
                return ConversionResult(
                    success=False,
                    error="pytesseract OCR produced no readable text",
                    complexity_report=report,
                    extraction_method="pytesseract-pdf",
                )
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="pytesseract-pdf",
                metadata={"engine": "pytesseract", "pages": len(md_parts)},
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"pytesseract PDF OCR failed: {e}",
                complexity_report=report,
                extraction_method="pytesseract-pdf",
            )

    def _convert_pdf_marker(self, file_path: str, report: FileComplexityReport, 
                            use_llm: bool = False, **options) -> ConversionResult:
        """Convert PDF using marker-pdf (best for OCR/scanned docs)."""
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.config.parser import ConfigParser
            
            config_dict = {"use_llm": use_llm} if use_llm else {}
            config_parser = ConfigParser(config_dict)
            models = create_model_dict()
            converter = PdfConverter(config=config_parser.generate_config_dict(), artifact_dict=models)
            rendered = converter(file_path)
            
            # Save images if requested
            images = []
            if options.get("extract_images") and hasattr(rendered, "images") and rendered.images:
                output_dir = options.get("image_output_dir", "./extracted_images")
                os.makedirs(output_dir, exist_ok=True)
                for name, img_data in rendered.images.items():
                    img_path = os.path.join(output_dir, name)
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    images.append({"name": name, "path": img_path, "size": len(img_data)})
            
            return ConversionResult(
                success=True,
                markdown=rendered.markdown,
                complexity_report=report,
                extraction_method="marker-pdf",
                metadata={"engine": "marker-pdf", "use_llm": use_llm},
                images_extracted=images,
                tables_extracted=[]
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"marker-pdf conversion failed: {e}",
                complexity_report=report,
                extraction_method="marker-pdf"
            )
    
    # === DOCX Converters ===
    
    def _convert_docx_native(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Convert DOCX using python-docx (native, preserves structure)."""
        try:
            from docx import Document
            from docx.oxml.ns import qn
            
            doc = Document(file_path)
            md_parts = []
            
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                
                style = para.style.name.lower() if para.style else ""
                
                # Handle headings
                if style.startswith("heading"):
                    level = 1
                    try:
                        level = int(style[-1]) if style[-1].isdigit() else 1
                    except:
                        pass
                    md_parts.append(f"{'#' * level} {para.text}")
                elif style == "title":
                    md_parts.append(f"# {para.text}")
                elif style == "subtitle":
                    md_parts.append(f"## {para.text}")
                else:
                    # Check for list formatting
                    text = para.text
                    if para.style and "list" in para.style.name.lower():
                        md_parts.append(f"- {text}")
                    else:
                        md_parts.append(text)
            
            # Handle tables
            tables = []
            for i, table in enumerate(doc.tables):
                table_md = self._table_to_markdown(table)
                tables.append({
                    "index": i,
                    "markdown": table_md,
                    "rows": len(table.rows),
                    "cols": len(table.columns)
                })
                md_parts.append(f"\n**Table {i+1}:**\n{table_md}\n")
            
            # Extract images
            images = []
            if options.get("extract_images"):
                images = self._extract_docx_images(doc, file_path)
                for img in images:
                    md_parts.append(f"\n![{img['name']}]({img['path']})\n")
            
            markdown = "\n\n".join(md_parts)
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="python-docx-native",
                metadata={"engine": "python-docx", "paragraphs": len(doc.paragraphs)},
                images_extracted=images,
                tables_extracted=tables
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"python-docx conversion failed: {e}",
                complexity_report=report,
                extraction_method="python-docx-native"
            )
    
    def _convert_docx_enhanced(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Enhanced DOCX conversion with better table/image handling."""
        opts = dict(options)
        opts.pop("extract_images", None)
        result = self._convert_docx_native(file_path, report, extract_images=True, **opts)
        if result.success and report.metadata.get("has_comments"):
            # Try to extract comments
            result.metadata["comments_extracted"] = self._extract_docx_comments(file_path)
        return result
    
    # === PPTX Converters ===
    
    def _convert_pptx_native(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Convert PPTX using python-pptx."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            md_parts = []
            tables = []
            images = []
            
            for slide_idx, slide in enumerate(prs.slides):
                md_parts.append(f"\n## Slide {slide_idx + 1}\n")
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                level = para.level + 1
                                md_parts.append(f"{'#' * min(level, 3)} {text}")
                    
                    if shape.has_table:
                        table_md = self._pptx_table_to_markdown(shape.table)
                        tables.append({
                            "slide": slide_idx + 1,
                            "markdown": table_md,
                            "rows": len(shape.table.rows),
                            "cols": len(shape.table.columns)
                        })
                        md_parts.append(f"\n**Table on Slide {slide_idx + 1}:**\n{table_md}\n")
                    
                    if hasattr(shape, 'image'):
                        # Save image
                        img = shape.image
                        ext = img.content_type.split('/')[-1]
                        img_name = f"slide{slide_idx+1}_shape{shape.shape_id}.{ext}"
                        img_path = os.path.join(options.get("image_output_dir", "./extracted_images"), img_name)
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                        with open(img_path, "wb") as f:
                            f.write(img.blob)
                        images.append({"name": img_name, "path": img_path, "slide": slide_idx + 1})
                        md_parts.append(f"\n![Slide {slide_idx+1} Image]({img_path})\n")
                
                # Speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        md_parts.append(f"\n**Speaker Notes:**\n{notes}\n")
            
            markdown = "\n".join(md_parts)
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="python-pptx-native",
                metadata={"engine": "python-pptx", "slides": len(prs.slides)},
                images_extracted=images,
                tables_extracted=tables
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"python-pptx conversion failed: {e}",
                complexity_report=report,
                extraction_method="python-pptx-native"
            )
    
    def _convert_pptx_enhanced(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Enhanced PPTX conversion."""
        opts = dict(options)
        opts.pop("extract_images", None)
        return self._convert_pptx_native(file_path, report, extract_images=True, **opts)
    
    # === XLSX Converters ===
    
    def _convert_xlsx_native(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Convert XLSX using openpyxl."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            md_parts = []
            tables = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                md_parts.append(f"\n## Sheet: {sheet_name}\n")
                
                # Convert to markdown table
                rows_data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows_data.append([str(cell) if cell is not None else "" for cell in row])
                
                if rows_data:
                    # Create markdown table
                    if len(rows_data) > 1:
                        header = rows_data[0]
                        separator = ["---"] * len(header)
                        table_md = "| " + " | ".join(header) + " |\n"
                        table_md += "| " + " | ".join(separator) + " |\n"
                        for row in rows_data[1:]:
                            table_md += "| " + " | ".join(row) + " |\n"
                    else:
                        table_md = "| " + " | ".join(rows_data[0]) + " |\n"
                    
                    tables.append({
                        "sheet": sheet_name,
                        "markdown": table_md,
                        "rows": len(rows_data),
                        "cols": len(rows_data[0]) if rows_data else 0
                    })
                    md_parts.append(table_md)
            
            markdown = "\n".join(md_parts)
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="openpyxl-native",
                metadata={"engine": "openpyxl", "sheets": len(wb.sheetnames)},
                images_extracted=[],
                tables_extracted=tables
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"openpyxl conversion failed: {e}",
                complexity_report=report,
                extraction_method="openpyxl-native"
            )
    
    def _convert_xlsx_enhanced(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Enhanced XLSX conversion with formula display."""
        try:
            import openpyxl
            
            # Load with formulas
            wb_formula = openpyxl.load_workbook(file_path, read_only=True, data_only=False)
            wb_data = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            md_parts = []
            tables = []
            
            for sheet_name in wb_formula.sheetnames:
                ws_f = wb_formula[sheet_name]
                ws_d = wb_data[sheet_name]
                
                md_parts.append(f"\n## Sheet: {sheet_name}\n")
                
                rows_data = []
                for row_idx, (row_f, row_d) in enumerate(zip(ws_f.iter_rows(values_only=False), ws_d.iter_rows(values_only=True))):
                    row_vals = []
                    for cell_f, cell_d in zip(row_f, row_d):
                        if cell_f.value is not None:
                            if isinstance(cell_f.value, str) and cell_f.value.startswith('='):
                                # Show formula
                                row_vals.append(f"`{cell_f.value}`")
                            else:
                                row_vals.append(str(cell_d) if cell_d is not None else "")
                        else:
                            row_vals.append("")
                    
                    if any(v for v in row_vals):
                        rows_data.append(row_vals)
                
                if rows_data:
                    header = rows_data[0]
                    separator = ["---"] * len(header)
                    table_md = "| " + " | ".join(header) + " |\n"
                    table_md += "| " + " | ".join(separator) + " |\n"
                    for row in rows_data[1:]:
                        table_md += "| " + " | ".join(row) + " |\n"
                    
                    tables.append({
                        "sheet": sheet_name,
                        "markdown": table_md,
                        "rows": len(rows_data),
                        "cols": len(rows_data[0]) if rows_data else 0
                    })
                    md_parts.append(table_md)
            
            markdown = "\n".join(md_parts)
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="openpyxl-enhanced",
                metadata={"engine": "openpyxl", "sheets": len(wb_formula.sheetnames), "formulas_shown": True},
                images_extracted=[],
                tables_extracted=tables
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"enhanced openpyxl conversion failed: {e}",
                complexity_report=report,
                extraction_method="openpyxl-enhanced"
            )
    
    # === Unstructured Converter ===
    
    def _convert_unstructured(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Convert using unstructured library."""
        try:
            from unstructured.partition.auto import partition
            from unstructured.documents.elements import Table, Image, Title, NarrativeText, ListItem
            
            elements = partition(filename=file_path)
            
            md_parts = []
            tables = []
            images = []
            table_idx = 0
            image_idx = 0
            
            for el in elements:
                if isinstance(el, Title):
                    md_parts.append(f"# {el.text}")
                elif isinstance(el, Table):
                    table_idx += 1
                    table_md = el.metadata.text_as_html if hasattr(el.metadata, 'text_as_html') else el.text
                    tables.append({
                        "index": table_idx,
                        "markdown": table_md,
                        "text": el.text
                    })
                    md_parts.append(f"\n**Table {table_idx}:**\n{table_md}\n")
                elif isinstance(el, Image):
                    image_idx += 1
                    # Save image if possible
                    if hasattr(el.metadata, 'image_base64') and el.metadata.image_base64:
                        img_data = base64.b64decode(el.metadata.image_base64)
                        img_name = f"unstructured_img_{image_idx}.png"
                        img_path = os.path.join(options.get("image_output_dir", "./extracted_images"), img_name)
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                        with open(img_path, "wb") as f:
                            f.write(img_data)
                        images.append({"name": img_name, "path": img_path, "base64": True})
                        md_parts.append(f"\n![Image {image_idx}]({img_path})\n")
                elif isinstance(el, (NarrativeText, ListItem)):
                    md_parts.append(el.text)
                elif el.text:
                    md_parts.append(el.text)
            
            markdown = "\n\n".join(md_parts)
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="unstructured",
                metadata={"engine": "unstructured", "elements": len(elements)},
                images_extracted=images,
                tables_extracted=tables
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"unstructured conversion failed: {e}",
                complexity_report=report,
                extraction_method="unstructured"
            )
    
    def _convert_unstructured_ocr(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Unstructured with OCR for scanned documents."""
        try:
            from unstructured.partition.auto import partition
            
            # Use OCR strategy
            elements = partition(
                filename=file_path,
                strategy="hi_res",  # High resolution for OCR
                ocr_languages=options.get("ocr_languages", ["eng"])
            )
            
            # Same processing as regular unstructured
            return self._process_unstructured_elements(elements, report, "unstructured-ocr", **options)
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"unstructured OCR conversion failed: {e}",
                complexity_report=report,
                extraction_method="unstructured-ocr"
            )
    
    def _process_unstructured_elements(self, elements, report, method, **options):
        """Process unstructured elements to markdown."""
        md_parts = []
        tables = []
        images = []
        table_idx = 0
        image_idx = 0
        
        from unstructured.documents.elements import Table, Image, Title, NarrativeText, ListItem
        
        for el in elements:
            if isinstance(el, Title):
                md_parts.append(f"# {el.text}")
            elif isinstance(el, Table):
                table_idx += 1
                table_md = el.metadata.text_as_html if hasattr(el.metadata, 'text_as_html') else el.text
                tables.append({"index": table_idx, "markdown": table_md, "text": el.text})
                md_parts.append(f"\n**Table {table_idx}:**\n{table_md}\n")
            elif isinstance(el, Image):
                image_idx += 1
                if hasattr(el.metadata, 'image_base64') and el.metadata.image_base64:
                    img_data = base64.b64decode(el.metadata.image_base64)
                    img_name = f"unstructured_img_{image_idx}.png"
                    img_path = os.path.join(options.get("image_output_dir", "./extracted_images"), img_name)
                    os.makedirs(os.path.dirname(img_path), exist_ok=True)
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    images.append({"name": img_name, "path": img_path})
                    md_parts.append(f"\n![Image {image_idx}]({img_path})\n")
            elif isinstance(el, (NarrativeText, ListItem)):
                md_parts.append(el.text)
            elif el.text:
                md_parts.append(el.text)
        
        markdown = "\n\n".join(md_parts)
        
        return ConversionResult(
            success=True,
            markdown=markdown,
            complexity_report=report,
            extraction_method=method,
            metadata={"engine": "unstructured", "elements": len(elements)},
            images_extracted=images,
            tables_extracted=tables
        )
    
    # === Image OCR ===
    
    def _convert_image_ocr(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Convert image to text using OCR."""
        if not self.backends["pytesseract"]:
            return ConversionResult(
                success=False,
                error="pytesseract not installed - cannot OCR images",
                complexity_report=report,
                extraction_method="pytesseract"
            )
        
        try:
            import pytesseract
            from PIL import Image
            
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, lang=options.get("ocr_lang", "eng"))
            
            markdown = f"# OCR Result: {Path(file_path).name}\n\n{text}"
            
            return ConversionResult(
                success=True,
                markdown=markdown,
                complexity_report=report,
                extraction_method="pytesseract-ocr",
                metadata={"engine": "pytesseract", "image_size": img.size, "lang": options.get("ocr_lang", "eng")},
                images_extracted=[{"name": Path(file_path).name, "path": file_path, "original": True}],
                tables_extracted=[]
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"OCR failed: {e}",
                complexity_report=report,
                extraction_method="pytesseract-ocr"
            )
    
    # === Text Files ===
    
    def _convert_text_file(self, file_path: str, report: FileComplexityReport, **options) -> ConversionResult:
        """Simple text/markdown file - just read it."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            
            return ConversionResult(
                success=True,
                markdown=content,
                complexity_report=report,
                extraction_method="direct-read",
                metadata={"engine": "direct-read", "encoding": "utf-8"},
                images_extracted=[],
                tables_extracted=[]
            )
        except Exception as e:
            return ConversionResult(
                success=False,
                error=f"Failed to read text file: {e}",
                complexity_report=report,
                extraction_method="direct-read"
            )
    
    # === Helper Methods ===
    
    def _extract_pdf_images(self, file_path: str, output_dir: str) -> List[Dict]:
        """Extract images from PDF using pymupdf."""
        images = []
        try:
            import pymupdf
            doc = pymupdf.open(file_path)
            os.makedirs(output_dir, exist_ok=True)
            
            for page_num, page in enumerate(doc):
                for img_idx, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    pix = pymupdf.Pixmap(doc, xref)
                    if pix.n >= 5:  # CMYK
                        pix = pymupdf.Pixmap(pymupdf.csRGB, pix)
                    img_name = f"page{page_num+1}_img{img_idx+1}.png"
                    img_path = os.path.join(output_dir, img_name)
                    pix.save(img_path)
                    images.append({
                        "name": img_name,
                        "path": img_path,
                        "page": page_num + 1,
                        "size": os.path.getsize(img_path)
                    })
        except Exception as e:
            pass
        return images
    
    def _extract_pdf_tables(self, file_path: str) -> List[Dict]:
        """Extract tables from PDF using pymupdf."""
        tables = []
        try:
            import pymupdf
            doc = pymupdf.open(file_path)
            
            for page_num, page in enumerate(doc):
                try:
                    found_tables = page.find_tables()
                    for table_idx, table in enumerate(found_tables.tables):
                        df = table.to_pandas()
                        table_md = df.to_markdown(index=False)
                        tables.append({
                            "page": page_num + 1,
                            "index": table_idx,
                            "markdown": table_md,
                            "rows": len(df),
                            "cols": len(df.columns)
                        })
                except:
                    pass
        except:
            pass
        return tables
    
    def _extract_docx_images(self, doc, file_path: str) -> List[Dict]:
        """Extract images from DOCX."""
        images = []
        try:
            output_dir = "./extracted_images"
            os.makedirs(output_dir, exist_ok=True)
            
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_data = rel.target_part.blob
                    content_type = rel.target_part.content_type
                    ext = content_type.split('/')[-1]
                    img_name = f"docx_img_{len(images)+1}.{ext}"
                    img_path = os.path.join(output_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    images.append({"name": img_name, "path": img_path, "size": len(img_data)})
        except:
            pass
        return images
    
    def _extract_docx_comments(self, file_path: str) -> List[Dict]:
        """Extract comments from DOCX."""
        comments = []
        try:
            from docx import Document
            doc = Document(file_path)
            if doc.part.comments_part:
                for comment in doc.part.comments_part.comments:
                    comments.append({
                        "id": comment.id,
                        "author": comment.author,
                        "text": comment.text,
                        "date": str(comment.date) if comment.date else None
                    })
        except:
            pass
        return comments
    
    def _table_to_markdown(self, table) -> str:
        """Convert python-docx table to markdown."""
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        
        if not rows:
            return ""
        
        header = rows[0]
        separator = ["---"] * len(header)
        md = "| " + " | ".join(header) + " |\n"
        md += "| " + " | ".join(separator) + " |\n"
        for row in rows[1:]:
            md += "| " + " | ".join(row) + " |\n"
        return md
    
    def _pptx_table_to_markdown(self, table) -> str:
        """Convert pptx table to markdown."""
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        
        if not rows:
            return ""
        
        header = rows[0]
        separator = ["---"] * len(header)
        md = "| " + " | ".join(header) + " |\n"
        md += "| " + " | ".join(separator) + " |\n"
        for row in rows[1:]:
            md += "| " + " | ".join(row) + " |\n"
        return md
    
    def _enhance_with_tables(self, file_path: str, result: ConversionResult, report: FileComplexityReport) -> Optional[ConversionResult]:
        """Add table extraction to existing result."""
        if not report.metadata.get("has_tables"):
            return None
        
        tables = self._extract_pdf_tables(file_path)
        if tables:
            result.tables_extracted = tables
            # Append tables to markdown
            table_md_parts = []
            for t in tables:
                table_md_parts.append(f"\n**Table on Page {t['page']}:**\n{t['markdown']}\n")
            result.markdown += "\n" + "\n".join(table_md_parts)
            result.metadata["tables_appended"] = True
            return result
        return None


def convert_to_markdown(file_path: str, **options) -> ConversionResult:
    """Convenience function for single file conversion."""
    converter = AnyDocMarkdownConverter()
    return converter.convert(file_path, **options)


def convert_bytes_to_markdown(file_bytes: bytes, filename: str, **options) -> ConversionResult:
    """Convenience function for bytes conversion."""
    converter = AnyDocMarkdownConverter()
    return converter.convert_bytes(file_bytes, filename, **options)


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python anydoc_markdown_converter.py <file_path> [--json] [--extract-images]")
        sys.exit(1)
    
    file_path = sys.argv[1]
    as_json = "--json" in sys.argv
    extract_images = "--extract-images" in sys.argv
    
    converter = AnyDocMarkdownConverter()
    result = converter.convert(file_path, extract_images=extract_images)
    
    if as_json:
        print(json.dumps(result.to_dict(), indent=2))
    else:
        if result.success:
            print(result.markdown)
        else:
            print(f"ERROR: {result.error}", file=sys.stderr)
            sys.exit(1)